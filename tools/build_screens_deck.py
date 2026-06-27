"""
Build a PowerPoint deck from the SAME slide source the web carousel uses
(marketing-aic/screens.js). Re-run any time the screenshots or copy change.

  python tools/build_screens_deck.py

Output: a 16:9 .pptx — one slide per screen (screenshot + heading + write-up),
plus a title slide. Editable in PowerPoint / Google Slides.
"""
import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

REPO = Path(__file__).resolve().parent.parent
MKT = REPO / "marketing-aic"
SRC = MKT / "screens.js"
OUT = Path.home() / "OneDrive" / "Desktop" / "AHQ screens" / "AccountsHQ-deck.pptx"

NAVY = RGBColor(0x0F, 0x16, 0x29)
GOLD = RGBColor(0xFC, 0xD3, 0x4D)
INK = RGBColor(0x0F, 0x17, 0x2A)
MUTE = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def load_slides():
    txt = SRC.read_text(encoding="utf-8")
    m = re.search(r"window\.AHQ_SCREENS\s*=\s*(\[.*\])\s*;", txt, re.S)
    return json.loads(m.group(1))


def add_textbox(slide, left, top, width, height, text, size, color, bold=False,
                anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return tb


def main():
    slides = load_slides()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W = prs.slide_width

    # ── Title slide ──
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(1, 0, 0, W, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    bg.shadow.inherit = False
    add_textbox(s, Inches(1), Inches(2.7), Inches(11.3), Inches(1.4),
                "See Accounts HQ in action", 44, GOLD, bold=True)
    add_textbox(s, Inches(1), Inches(4.1), Inches(11.3), Inches(1.0),
                "Real screens, real data — what your clients' day looks like.",
                20, RGBColor(0xCB, 0xD5, 0xE1))

    # ── One slide per screen: screenshot left, write-up right ──
    for sc in slides:
        s = prs.slides.add_slide(blank)
        img = MKT / sc["img"]
        # screenshot — left, 7.5" wide (aspect kept), vertically centred-ish
        if img.exists():
            pic = s.shapes.add_picture(str(img), Inches(0.45), Inches(1.7),
                                       width=Inches(7.6))
        add_textbox(s, Inches(8.35), Inches(0.9), Inches(4.55), Inches(1.6),
                    sc["title"], 24, INK, bold=True)
        add_textbox(s, Inches(8.35), Inches(2.7), Inches(4.55), Inches(4.4),
                    sc["body"], 14, MUTE)
        # thin gold rule under the title
        ln = s.shapes.add_shape(1, Inches(8.35), Inches(2.45), Inches(1.2), Pt(3))
        ln.fill.solid(); ln.fill.fore_color.rgb = GOLD; ln.line.fill.background()
        ln.shadow.inherit = False

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {len(slides)} content slides -> {OUT}")


if __name__ == "__main__":
    main()
